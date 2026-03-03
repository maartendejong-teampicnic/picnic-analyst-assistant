#!/usr/bin/env python3
"""
Slides CLI tool for Claude Code.
Creates .pptx files from a JSON deck definition.

Usage:
  slides_tool.py <deck.json>
  slides_tool.py --json '<json_string>'
  echo '<json>' | slides_tool.py -
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

# Maps user-facing layout keys to the Picnic template layout names
LAYOUT_MAP = {
    "start":      "Start",
    "agenda":     "Agenda",
    "title-only": "Title Only",
    "bullets":    "Bullets",
    "charts":     "Charts",
    "blank":      "Blank",
    "end":        "End",
}

# For each template layout, which placeholder index serves each logical role.
# Based on inspection of picnic_default_ppt_theme.pptx and analysis of real decks:
#
#   ph_idx  0  = Rectangle 5 — the thin colored top strip; in content slides this
#                holds the MAIN TITLE / CONCLUSION (not a visual element — it's text)
#   ph_idx 10  = Text below the strip — the SUBTITLE (metric definition, scope,
#                secondary point) OR the main title in Start/Agenda layouts
#   ph_idx 11  = Second line of the Start layout (subtitle on opening slide)
#   ph_idx 23  = Footer — actively used for data scope notes, dataset descriptions
#   ph_idx 24  = Body / bullets (Bullets layout) OR sidebar box 1 (Charts)
#   ph_idx 25  = sidebar box 2 (Charts)
#   ph_idx 26  = sidebar box 3 (Charts) / agenda item 1 text (Agenda)
#   ph_idx 32/34/36/38/40/42 = agenda item numbers 1–6
#   ph_idx 33/35/37/39/41    = agenda item text 2–6
PLACEHOLDER_ROLES = {
    "Start": {
        "title":    10,   # Main title (positioned at bottom of slide)
        "subtitle": 11,   # Subtitle below main title
    },
    "Agenda": {
        "title":       10,
        "item_text":   [26, 33, 35, 37, 39, 41],   # up to 6 agenda items
        "item_number": [32, 34, 36, 38, 40, 42],   # auto-numbered
    },
    "Title Only": {
        "title":    0,    # Conclusion / main finding (thin strip at top)
        "subtitle": 10,   # Metric definition, scope, or secondary point
        "footer":   23,   # Data source / dataset scope note
    },
    "Bullets": {
        "title":    0,    # Conclusion
        "subtitle": 10,   # Secondary text
        "body":     24,   # Bullet points
        "footer":   23,   # Data note
    },
    "Charts": {
        "title":    0,    # Conclusion
        "subtitle": 10,   # Secondary text
        "footer":   23,   # Data note
        "sidebar":  [24, 25, 26],   # Right-side annotation text boxes
    },
    "Blank": {
        "footer":   23,
    },
    "End": {},
}


def find_layout(prs, layout_key: str):
    """Return the slide layout matching the given key. Raises on unknown key."""
    target_name = LAYOUT_MAP.get(layout_key.lower())
    if target_name is None:
        raise ValueError(
            f"Unknown layout {layout_key!r}. Supported: {list(LAYOUT_MAP.keys())}"
        )
    for layout in prs.slide_layouts:
        if layout.name == target_name:
            return layout
    available = [l.name for l in prs.slide_layouts]
    raise ValueError(
        f"Layout {target_name!r} not found in this template. "
        f"Template has: {available}"
    )


def set_text(slide, ph_idx: int, text: str) -> bool:
    """Set plain text in a placeholder by index. Returns True if found."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            tf.clear()
            tf.text = str(text)
            return True
    return False


def set_bullets(slide, ph_idx: int, bullets: list) -> bool:
    """Set multi-paragraph bullet content in a placeholder."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            tf.clear()
            if not bullets:
                return True
            # First paragraph (tf.clear() leaves one empty paragraph)
            first = bullets[0]
            p0 = tf.paragraphs[0]
            if isinstance(first, dict):
                p0.text  = first.get("text", "")
                p0.level = first.get("level", 0)
            else:
                p0.text = str(first)
            # Additional paragraphs
            for item in bullets[1:]:
                p = tf.add_paragraph()
                if isinstance(item, dict):
                    p.text  = item.get("text", "")
                    p.level = item.get("level", 0)
                else:
                    p.text = str(item)
            return True
    return False


_NOTES_PH_XML = (
    '<p:sp '
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<p:nvSpPr>"
    '<p:cNvPr id="3" name="Notes Placeholder 2"/>'
    "<p:cNvSpPr><a:spLocks noGrp=\"1\"/></p:cNvSpPr>"
    "<p:nvPr><p:ph type=\"body\" idx=\"1\"/></p:nvPr>"
    "</p:nvSpPr>"
    "<p:spPr>"
    "<a:xfrm><a:off x=\"457200\" y=\"3600450\"/>"
    "<a:ext cx=\"8229600\" cy=\"3547650\"/></a:xfrm>"
    "</p:spPr>"
    "<p:txBody><a:bodyPr/><a:lstStyle/>"
    "<a:p><a:endParaRPr lang=\"en-US\" dirty=\"0\"/></a:p>"
    "</p:txBody>"
    "</p:sp>"
)


def set_notes(slide, text: str):
    """Set speaker notes text, injecting a notes placeholder if the template lacks one."""
    if not text:
        return
    from lxml import etree

    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf is None:
        # Template has no body placeholder on notes slides — add one
        notes_slide._element.cSld.spTree.append(etree.fromstring(_NOTES_PH_XML))
        tf = notes_slide.notes_text_frame
    if tf is not None:
        tf.text = str(text)


def add_slide(prs, slide_def: dict):
    """Add one slide to the presentation based on the slide definition dict."""
    layout_key = slide_def.get("layout", "bullets")
    layout = find_layout(prs, layout_key)
    layout_name = layout.name
    slide = prs.slides.add_slide(layout)
    roles = PLACEHOLDER_ROLES.get(layout_name, {})

    # Title
    title = slide_def.get("title", "")
    if title and "title" in roles:
        set_text(slide, roles["title"], title)

    # Subtitle (Start, Title Only, Bullets, Charts)
    subtitle = slide_def.get("subtitle", "")
    if subtitle and "subtitle" in roles:
        set_text(slide, roles["subtitle"], subtitle)

    # Footer / data scope note (Title Only, Bullets, Charts, Blank)
    footer = slide_def.get("footer", "")
    if footer and "footer" in roles:
        set_text(slide, roles["footer"], footer)

    # Body / bullets (Bullets layout)
    body = slide_def.get("body", [])
    if body and "body" in roles:
        if isinstance(body, str):
            body = [body]
        set_bullets(slide, roles["body"], body)

    # Agenda items (Agenda layout)
    items = slide_def.get("items", [])
    if items and "item_text" in roles:
        for text_idx, num_idx, item in zip(
            roles["item_text"], roles["item_number"], items
        ):
            set_text(slide, text_idx, item)
            set_text(slide, num_idx, str(items.index(item) + 1))

    # Sidebar annotation boxes (Charts layout)
    sidebar = slide_def.get("sidebar", [])
    if sidebar and "sidebar" in roles:
        for ph_idx, content in zip(roles["sidebar"], sidebar):
            if isinstance(content, list):
                set_bullets(slide, ph_idx, content)
            else:
                set_text(slide, ph_idx, str(content))

    # Speaker notes
    set_notes(slide, slide_def.get("notes", ""))
    return slide


def build_deck(deck_def: dict) -> str:
    """Build the .pptx file from the deck definition. Returns output path."""
    output_path = deck_def.get("output")
    if not output_path:
        raise ValueError("Deck definition must include an 'output' path.")

    template_path = deck_def.get("template")
    if template_path:
        prs = Presentation(str(Path(template_path).expanduser()))
        # Drop each existing slide: remove both the relationship and the sldId element
        # so the original XML parts don't end up as duplicates in the saved zip.
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        sld_id_lst = prs.slides._sldIdLst
        for sld_id in list(sld_id_lst):
            rId = sld_id.get(f"{{{r_ns}}}id")
            if rId:
                prs.part.drop_rel(rId)
            sld_id_lst.remove(sld_id)
    else:
        prs = Presentation()

    for slide_def in deck_def.get("slides", []):
        add_slide(prs, slide_def)

    output = Path(output_path).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    n = len(prs.slides)
    print(f"Saved: {output}  ({n} slide{'s' if n != 1 else ''})")
    return str(output)


def main():
    parser = argparse.ArgumentParser(
        description="Create PowerPoint .pptx files from a JSON deck definition"
    )
    parser.add_argument(
        "input", nargs="?",
        help="Path to JSON deck file, or '-' to read from stdin"
    )
    parser.add_argument("--json", metavar="JSON", help="Inline JSON deck definition")
    args = parser.parse_args()

    if args.json:
        deck_def = json.loads(args.json)
    elif args.input == "-" or (not args.input and not sys.stdin.isatty()):
        deck_def = json.loads(sys.stdin.read())
    elif args.input:
        deck_def = json.loads(Path(args.input).expanduser().read_text())
    else:
        parser.print_help()
        sys.exit(1)

    build_deck(deck_def)


if __name__ == "__main__":
    main()
