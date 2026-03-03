#!/usr/bin/env python3
"""
Think-cell automation tool for Claude Code.

Generates a .ppttc data file and calls ppttc.exe to produce a .pptx with
real think-cell charts and text fields filled from structured JSON data.
Optionally post-processes the output with python-pptx to fill slide text
placeholders (title, subtitle, footer) that are not think-cell elements.

Usage:
  thinkcell_tool.py <deck.json>
  thinkcell_tool.py --json '<json_string>'
  echo '<json>' | thinkcell_tool.py -

Input JSON format:
{
  "template": "/path/to/template.pptx",   // .pptx with named think-cell elements
  "output":   "/path/to/output.pptx",
  "charts": [
    {
      "name": "conversion_bar",            // name assigned in think-cell
      "categories": ["Control", "Test"],
      "series": [
        {"name": "Conversion (%)", "values": [81.1, 83.1]}
      ]
    }
  ],
  "textfields": [
    {"name": "test_period", "value": "CW05–08 2026"}
  ],
  "slides": [
    {
      "slide_index": 3,      // 0-based slide number
      "title":    "...",     // ph[0]  — main conclusion
      "subtitle": "...",     // ph[10] — metric definition
      "footer":   "..."      // ph[23] — data scope note
    }
  ]
}

Naming think-cell elements in PowerPoint:
  Right-click a think-cell chart or text field → "Name" (or "Add Name")
  The name you assign here must match the "name" field in the JSON above.
"""

import argparse
import json
import subprocess
import sys
from pathlib import Path

PPTTC_EXE = r"C:\Program Files (x86)\think-cell\ppttc.exe"


def wsl_to_win(path: str) -> str:
    """Convert a WSL/Linux path to a Windows drive path (C:\\...) via wslpath."""
    result = subprocess.run(
        ["wslpath", "-w", str(path)],
        capture_output=True, text=True, check=True,
    )
    return result.stdout.strip()


def win_temp_dir_wsl() -> Path:
    """Return the Windows %TEMP% directory as a WSL-accessible path under /mnt/."""
    win_temp = subprocess.run(
        ["powershell.exe", "-Command", "[System.IO.Path]::GetTempPath()"],
        capture_output=True, text=True, check=True,
    ).stdout.strip().rstrip("\\")
    wsl_temp = subprocess.run(
        ["wslpath", win_temp],
        capture_output=True, text=True, check=True,
    ).stdout.strip()
    return Path(wsl_temp)


def _cell(val):
    """Wrap a value in the ppttc cell object format required by think-cell."""
    if val is None:
        return None
    if isinstance(val, (int, float)):
        return {"number": val}
    return {"string": str(val)}


def build_ppttc_data(deck_def: dict) -> list:
    """
    Convert the deck definition into the think-cell data array.

    Two table orientations depending on chart type:

    "bar" (default) — wide format, categories as columns:
      Row 0: [null, {"string": cat1}, {"string": cat2}, ...]
      Row 1+: [{"string": series_name}, {"number": val1}, ...]

    "line" — tall format, categories as rows (required for line/scatter charts):
      Row 0: [null, {"string": series1_name}, ...]   ← column headers
      Row 1: [null, null, ...]                        ← empty 2nd header row
      Row 2+: [{"string": cat}, {"number": val}, ...]  ← one row per category

    Specify "type": "line" in the chart definition to use tall format.
    """
    data = []

    for chart in deck_def.get("charts", []):
        categories  = chart.get("categories", [])
        series      = chart.get("series", [])
        chart_type  = chart.get("type", "bar")

        if chart_type == "line":
            # Line chart format: categories as columns, empty row 2, data in row 3
            #   Row 0: [null, cat1, cat2, ...]
            #   Row 1: [null, null, ...]          ← required empty separator
            #   Row 2: [series_name, val1, val2, ...]
            n = len(categories)
            table = (
                [[None] + [_cell(c) for c in categories]]
                + [[None] + [None] * n]
                + [[_cell(s["name"])] + [_cell(v) for v in s["values"]] for s in series]
            )
        else:
            # Bar/column chart format: categories as columns, series as rows
            table = (
                [[None] + [_cell(c) for c in categories]]
                + [[_cell(s["name"])] + [_cell(v) for v in s["values"]] for s in series]
            )
        data.append({"name": chart["name"], "table": table})

    for tf in deck_def.get("textfields", []):
        data.append({"name": tf["name"], "string": str(tf["value"])})

    return data


def run(deck_def: dict) -> str:
    """Build the deck and return the output path."""
    template_path = deck_def.get("template")
    output_path   = deck_def.get("output")
    if not template_path:
        raise ValueError("Deck definition must include a 'template' path.")
    if not output_path:
        raise ValueError("Deck definition must include an 'output' path.")

    template = Path(template_path).expanduser().resolve()
    output   = Path(output_path).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)

    if not template.exists():
        raise FileNotFoundError(f"Template not found: {template}")

    # ppttc.exe requires real Windows drive paths (C:\...), not WSL UNC paths.
    # Stage template + .ppttc in the Windows %TEMP% folder, run there, then
    # copy the result to the requested output path.
    import shutil, uuid
    run_id   = uuid.uuid4().hex[:8]
    tmp      = win_temp_dir_wsl() / f"tc_{run_id}"
    tmp.mkdir(parents=True, exist_ok=True)

    tmp_template = tmp / ("template" + template.suffix)
    tmp_ppttc    = tmp / "deck.ppttc"
    tmp_output   = tmp / "output.pptx"

    shutil.copy2(template, tmp_template)
    template_win = wsl_to_win(str(tmp_template))
    ppttc_win    = wsl_to_win(str(tmp_ppttc))
    output_win   = wsl_to_win(str(tmp_output))

    ppttc_content = [{"template": template_win, "data": build_ppttc_data(deck_def)}]
    tmp_ppttc.write_text(
        json.dumps(ppttc_content, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    try:
        print(f"Running ppttc.exe → {output.name} ...")
        result = subprocess.run(
            [
                "powershell.exe", "-Command",
                f'& "{PPTTC_EXE}" "{ppttc_win}" -o "{output_win}"',
            ],
            capture_output=True,
            text=True,
        )
    finally:
        # Copy result out before cleaning up, even on failure
        if tmp_output.exists():
            shutil.copy2(tmp_output, output)
        shutil.rmtree(tmp, ignore_errors=True)

    if result.returncode != 0:
        stderr = result.stderr.strip() or result.stdout.strip()
        raise RuntimeError(
            f"ppttc.exe failed (exit {result.returncode}):\n{stderr}"
        )

    if result.stdout.strip():
        print(result.stdout.strip())

    print(f"Saved: {output}")

    # Post-process: fill slide text placeholders via python-pptx
    slides_text = deck_def.get("slides", [])
    if slides_text:
        _fill_slide_text(output, slides_text)

    return str(output)


# Placeholder index mapping (matches Picnic template layout "Title Only")
_PH_ROLES = {"title": 0, "subtitle": 10, "footer": 23}


def _fill_slide_text(pptx_path: Path, slides_text: list) -> None:
    """Fill title / subtitle / footer placeholders on specified slides."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    for slide_spec in slides_text:
        idx = slide_spec.get("slide_index")
        if idx is None or idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        for field, ph_idx in _PH_ROLES.items():
            text = slide_spec.get(field)
            if text and ph_idx in ph_map:
                ph_map[ph_idx].text_frame.text = text

    prs.save(str(pptx_path))
    print(f"Slide text filled on {len(slides_text)} slide(s).")


def main():
    parser = argparse.ArgumentParser(
        description="Fill a think-cell template with JSON data and produce a .pptx"
    )
    parser.add_argument(
        "input", nargs="?",
        help="Path to JSON deck file, or '-' to read from stdin",
    )
    parser.add_argument("--json", metavar="JSON", help="Inline JSON deck definition")
    args = parser.parse_args()

    if args.json:
        deck_def = json.loads(args.json)
    elif args.input == "-" or (not args.input and not sys.stdin.isatty()):
        deck_def = json.loads(sys.stdin.read())
    elif args.input:
        deck_def = json.loads(Path(args.input).expanduser().read_text())
    else:
        parser.print_help()
        sys.exit(1)

    run(deck_def)


if __name__ == "__main__":
    main()
